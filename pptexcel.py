from flask import Flask, render_template, request, send_file
from pptx import Presentation
import re
import os
from openpyxl import load_workbook
from openpyxl.utils import get_column_letter

app = Flask(__name__)

def copy_row_format(ws, source_row, target_row):
    """Copy the format of the source_row to target_row."""
    for col_num in range(1, ws.max_column + 1):
        col_letter = get_column_letter(col_num)
        source_cell = ws[f'{col_letter}{source_row}']
        target_cell = ws[f'{col_letter}{target_row}']

        # Copy font
        if source_cell.font:
            target_cell.font = source_cell.font.copy()

        # Copy border
        if source_cell.border:
            target_cell.border = source_cell.border.copy()

        # Copy fill
        if source_cell.fill:
            target_cell.fill = source_cell.fill.copy()

        # Copy number format
        if source_cell.number_format:
            target_cell.number_format = source_cell.number_format

        # Copy protection
        if source_cell.protection:
            target_cell.protection = source_cell.protection.copy()

        # Copy alignment
        if source_cell.alignment:
            target_cell.alignment = source_cell.alignment.copy()


@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        ppt_files = request.files.getlist('ppt_files')
        formatted_excel = request.files['formatted_excel']
        company_name = request.form['company_name']

        # Load the formatted Excel file using openpyxl
        wb = load_workbook(formatted_excel)
        ws = wb.active  # Assuming you want to modify the first sheet

        for ppt_file in ppt_files:
            # Save the uploaded file temporarily
            ppt_path = os.path.join('temp', ppt_file.filename)
            ppt_file.save(ppt_path)

            # Load the PowerPoint presentation
            presentation = Presentation(ppt_path)

            # Extracting data from slide 3 (Trainer’s Name, Workshop Title, Date, and Location)
            slide_3 = presentation.slides[2]
            slide_3_data = {}
            for shape in slide_3.shapes:
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if "Trainer" in text:
                        match = re.search(r"Trainer’s Name: ([\w\s\.]+)(?=\s+Workshop Title:)", text)
                        match2 = re.search(r"Workshop Title:\s*(.+)$", text)
                        if match2:
                            workshop_title = match2.group(1).strip()
                            slide_3_data['Workshop Title'] = workshop_title
                        if match:
                            trainer_name = match.group(1).strip()
                            trainer_title = trainer_name.split()[0]  # Extract first string as Trainer Title
                            trainer_name = trainer_name[len(trainer_title):].strip()
                            slide_3_data['Trainer’s Name'] = trainer_name
                            slide_3_data['Trainer Title'] = trainer_title
                    elif "Date" in text:
                        date_match = re.search(r"Date:\s*(.*?)\s*Location:", text)
                        location_match = re.search(r"Location:\s*(.+)$", text)
                        if location_match:
                            workshop_location = location_match.group(1).strip()
                            slide_3_data['Location'] = workshop_location
                        if date_match:
                            workshop_date = date_match.group(1).strip()
                            slide_3_data['Date'] = workshop_date
            slide_7 = presentation.slides[6]
            y_axis_values = []
            for shape in slide_7.shapes:
                if shape.has_chart:
                    chart = shape.chart
                    for series in chart.series:
                        y_axis_values.extend(series.values)
            # Extracting data from slide 8 (Trainer percentage)
            slide_8 = presentation.slides[7]
            trainer_percentage = None
            for shape in slide_8.shapes:
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if "Trainer" in text:
                        break
                    if "%" in text:
                        trainer_percentage = text

            # Prepare the data to insert into the formatted Excel file
            new_row = [
                company_name,
                slide_3_data.get('Trainer Title', ''),
                slide_3_data.get('Trainer’s Name', ''),
                slide_3_data.get('Workshop Title', ''),
                slide_3_data.get('Date', ''),
                slide_3_data.get('Location', ''),
                trainer_percentage,
                str(y_axis_values[0] * 100) + '%',  # Instructor was knowledgeable and competent about the topic.
                str(y_axis_values[1] * 100) + '%',  # Instructor was organized, grabbed attention and built rapport.
                str(y_axis_values[2] * 100) + '%',  # Instructor allowed time for interaction and participation.
                str(y_axis_values[3] * 100) + '%'   # Instructor created a positive and safe learning environment.
            ]
            trainer_name_sheet_found = False
            for sheet_name in wb.sheetnames:
                sheet_name_trimmed = sheet_name.strip()  # Remove any leading or trailing spaces
                if sheet_name_trimmed.lower() == trainer_name.lower():
                    ws_trainer = wb[sheet_name]
                    trainer_name_sheet_found = True
                    break
            if trainer_name_sheet_found:
                # Find the next empty row in the trainer's sheet
                next_row = ws_trainer.max_row + 1

                # Copy the format of the previous row to the new row
                if next_row > 2:  # Ensure there is a row to copy from
                    copy_row_format(ws_trainer, next_row - 1, next_row)

                # Insert the data into the new row
                for col_num, value in enumerate(new_row, start=1):
                    ws_trainer.cell(row=next_row, column=col_num, value=value)

            # Find the next empty row in the Excel sheet
            next_row = ws.max_row + 1

            # Copy the format of the previous row to the new row
            if next_row > 2:  # Ensure there is a row to copy from
                copy_row_format(ws, next_row - 1, next_row)

            # Insert the data into the new row
            for col_num, value in enumerate(new_row, start=1):
                ws.cell(row=next_row, column=col_num, value=value)

        # Save the updated formatted Excel file
        formatted_excel_path = os.path.join('temp', 'Updated_Formatted_Excel.xlsx')
        wb.save(formatted_excel_path)

        # Send the updated file to the user for download
        return send_file(formatted_excel_path, as_attachment=True, download_name='Trainers Analysis.xlsx')

    return render_template('analysis.html')

if __name__ == '__main__':
    if not os.path.exists('temp'):
        os.makedirs('temp')
    app.run(debug=True)
